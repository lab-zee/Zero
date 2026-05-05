"""
File content extraction utilities for reading text from various file formats.
"""
import os
from pathlib import Path
from typing import Optional
from .storage import get_file

def extract_text_from_file(file_path: str, content_type: Optional[str] = None, original_filename: str = "") -> tuple[str, Optional[str]]:
    """
    Extract text content from a file based on its type.
    Returns (extracted_text, error_message) tuple.
    If successful, error_message is None. If failed, extracted_text is empty string.
    """
    file_ext = Path(original_filename or file_path).suffix.lower()
    
    try:
        # Check if file exists
        import os
        if not os.path.exists(file_path):
            return "", f"File not found at path: {file_path}"
        
        file_content = get_file(file_path)
        
        if not file_content or len(file_content) == 0:
            return "", f"File is empty: {original_filename}"
        
        # Text files
        if file_ext in ['.txt', '.csv']:
            try:
                return file_content.decode('utf-8'), None
            except UnicodeDecodeError:
                try:
                    return file_content.decode('latin-1'), None
                except Exception as e:
                    return "", f"Could not decode {file_ext} file: {str(e)}"
        
        # PDF files - try pypdf first (better), fallback to PyPDF2
        elif file_ext == '.pdf':
            try:
                # Try pypdf first (newer, better library)
                try:
                    from pypdf import PdfReader
                except ImportError:
                    # Fallback to PyPDF2
                    import PyPDF2
                    PdfReader = PyPDF2.PdfReader
                
                from io import BytesIO
                pdf_reader = PdfReader(BytesIO(file_content))
                text_parts = []
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text.strip():  # Only add non-empty pages
                        text_parts.append(text)
                
                if not text_parts:
                    return "", "PDF appears to be empty or contains only images. Could not extract text."
                
                return "\n".join(text_parts), None
            except ImportError:
                return "", "PDF library not installed. Install with: pip install pypdf or pip install PyPDF2"
            except Exception as e:
                return "", f"Error extracting PDF: {str(e)}. The PDF may be encrypted, corrupted, or contain only images."
        
        # Excel files
        elif file_ext in ['.xls', '.xlsx']:
            try:
                import openpyxl
                from io import BytesIO
                workbook = openpyxl.load_workbook(BytesIO(file_content), data_only=True)
                text_parts = []
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_parts.append(f"\n--- Sheet: {sheet_name} ---\n")
                    for row in sheet.iter_rows(values_only=True):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        text_parts.append("\t".join(row_data))
                return "\n".join(text_parts), None
            except ImportError:
                return "", "openpyxl library not installed. Install with: pip install openpyxl"
            except Exception as e:
                return "", f"Error extracting Excel: {str(e)}"
        
        # Word documents
        elif file_ext in ['.doc', '.docx']:
            try:
                from docx import Document
                from io import BytesIO
                doc = Document(BytesIO(file_content))
                text_parts = []
                for paragraph in doc.paragraphs:
                    text_parts.append(paragraph.text)
                return "\n".join(text_parts), None
            except ImportError:
                return "", "python-docx library not installed. Install with: pip install python-docx"
            except Exception as e:
                return "", f"Error extracting Word document: {str(e)}"
        
        # PowerPoint
        elif file_ext in ['.ppt', '.pptx']:
            try:
                from pptx import Presentation
                from io import BytesIO
                prs = Presentation(BytesIO(file_content))
                text_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    text_parts.append(f"\n--- Slide {slide_num} ---\n")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                return "\n".join(text_parts), None
            except ImportError:
                return "", "python-pptx library not installed. Install with: pip install python-pptx"
            except Exception as e:
                return "", f"Error extracting PowerPoint: {str(e)}"
        
        else:
            return "", f"File type {file_ext} is not yet supported for text extraction"
    
    except Exception as e:
        return "", f"Error reading file: {str(e)}"

def extract_files_content(db_files: list) -> tuple[str, list[str]]:
    """
    Extract text content from multiple files and combine them.
    Returns (combined_content, error_messages) tuple.
    error_messages is a list of error strings for files that failed to extract.
    """
    if not db_files:
        return "", []
    
    content_parts = []
    error_messages = []
    
    for db_file in db_files:
        print(f"Extracting content from file: {db_file.original_filename} (path: {db_file.file_path})")
        file_content, error = extract_text_from_file(
            db_file.file_path,
            db_file.content_type,
            db_file.original_filename
        )
        
        if error:
            error_msg = f"Error extracting {db_file.original_filename}: {error}"
            print(f"WARNING: {error_msg}")
            error_messages.append(error_msg)
            # Still include the filename so LLM knows a file was attempted
            content_parts.append(f"\n--- File: {db_file.original_filename} ---\n[EXTRACTION FAILED: {error}]\n")
        elif file_content:
            print(f"Successfully extracted {len(file_content)} characters from {db_file.original_filename}")
            content_parts.append(f"\n--- File: {db_file.original_filename} ---\n{file_content}\n")
        else:
            error_msg = f"No content extracted from {db_file.original_filename}"
            print(f"WARNING: {error_msg}")
            error_messages.append(error_msg)
    
    result = "\n".join(content_parts)
    print(f"Total extracted content length: {len(result)} characters")
    if error_messages:
        print(f"Extraction errors: {error_messages}")
    
    return result, error_messages

